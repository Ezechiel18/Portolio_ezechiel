#!/usr/bin/env python3
# =============================================================================
#  main.py — Point d'entrée unique — Qualification SIG SCE — V5
#
#  CHANGEMENTS V5 :
#    - Dédoublonnage des exemples en section 2.4 (combinaison couche+champ+valeur)
#    - Génération PDF via Chrome headless (1 page, fidèle au HTML)
#    - Génération PowerPoint éditable via python-pptx
#    - Commentaires ligne par ligne sur tout le fichier
#
#  FICHIERS À MODIFIER SELON LE PROJET :
#    - mapping.py  : noms des champs client (SEULE MODIFICATION NÉCESSAIRE)
#    - config.py   : chemins, nom projet, commune, logo SCE
#
#  USAGE :
#    python main.py
#    python main.py --troncon "chemin/Troncon.shp" --noeud "chemin/Noeud.shp"
# =============================================================================

import os           # gestion des chemins système
import sys          # manipulation du chemin Python pour les imports
import argparse     # gestion des arguments ligne de commande
import locale       # pour la date en français
import subprocess   # pour appeler Chrome en ligne de commande
import shutil       # pour vérifier si Chrome est installé
from datetime import datetime   # date générée automatiquement
from pathlib import Path        # gestion moderne des chemins
from base64 import b64encode    # encodage des images pour le PDF

import geopandas as gpd                               # lecture fichiers SIG
from jinja2 import Environment, FileSystemLoader      # moteur de templates HTML
from xhtml2pdf import pisa                            # PDF de secours si Chrome absent

# --- Chemins du projet ---
ROOT    = Path(__file__).parent          # dossier racine du projet
OUT_DIR = ROOT / "output"               # dossier de sortie principal
CHT_DIR = OUT_DIR / "charts"            # sous-dossier graphiques PNG
MAP_DIR = OUT_DIR / "maps"              # sous-dossier cartes PNG
TPL_DIR = ROOT / "templates"           # dossier du template HTML

# Ajout du dossier racine au chemin Python pour les imports relatifs
sys.path.insert(0, str(ROOT))

# Import de la configuration et du mapping
from config  import (
    TRONCON_FILE, NOEUD_FILE,           # chemins des données client
    PROJET_NOM, PROJET_COMMUNE,         # infos projet pour le rapport
    PROJET_VERSION, LOGO_SCE_PATH,      # version et logo
    POIDS_ATTR, POIDS_TOPO,             # pondérations IFQ
    ACTIONS_CORRECTIVES                 # dictionnaire des actions correctives par problème
)
from mapping import CHAMPS_TRONCON, CHAMPS_NOEUD   # correspondance champs client

# Import des modules d'analyse
from analysis.quality_attributaire import (
    run_quality_attributaire, _preparer_icgp, arrondi_pct, _preparer_synthese_4_4
)
from analysis.quality_spatiale      import run_quality_spatiale

# Import des modules de visualisation
from viz.charts import (
    jauge_demi_cercle, donut,
    barres_horizontales, couleur_score
)
from viz.generate_pptx import generer_pptx
from viz.maps import (
    carte_noeuds_connectivite,
    carte_troncons_delimitation
)


# =============================================================================
#  ARGUMENTS LIGNE DE COMMANDE
#  Permettent de surcharger les chemins de config.py sans modifier le code.
# =============================================================================
def parse_args():
    """
    Définit les arguments optionnels.
    Si non fournis, les valeurs de config.py sont utilisées.
    """
    p = argparse.ArgumentParser(description="Qualification SIG — SCE V8")
    p.add_argument("--troncon", default=TRONCON_FILE,
                   help="Chemin vers le fichier tronçon du client")
    p.add_argument("--noeud",   default=NOEUD_FILE,
                   help="Chemin vers le fichier nœud du client")
    p.add_argument("--out",     default=str(OUT_DIR),
                   help="Dossier de sortie pour le rapport")
    return p.parse_args()


# =============================================================================
#  CHARGEMENT D'UNE COUCHE SIG
# =============================================================================
def charger_couche(chemin, nom):
    """
    Charge un fichier SIG avec GeoPandas.
    Arrête le script avec un message clair si le fichier est introuvable.
    Supporte : .shp, .gpkg, et tout format reconnu par GDAL/GeoPandas.
    """
    chemin = Path(chemin)
    if not chemin.exists():
        # Message d'erreur explicite pour aider l'utilisateur
        print(f"\n  [ERREUR] Fichier introuvable : {chemin}")
        print("  → Vérifier DATA_DIR dans config.py ou utiliser --troncon / --noeud")
        sys.exit(1)

    print(f"  Chargement {nom} … ", end="", flush=True)
    gdf = gpd.read_file(chemin)   # lecture du fichier SIG
    print(f"{len(gdf)} objets  |  CRS : {gdf.crs}")
    return gdf


# =============================================================================
#  VALIDATION STRICTE DES CHAMPS DÉFINIS DANS MAPPING.PY
# =============================================================================
def valider_champs_mapping(gdf, champs_config, nom_couche):
    """
    Vérifie que chaque "nom_client" renseigné dans mapping.py (CHAMPS_TRONCON
    ou CHAMPS_NOEUD) correspond bien à une colonne réelle de la couche.

    Distinction stricte entre deux cas qui ne doivent jamais être confondus :
      - nom_client = None      → choix volontaire, le client n'a pas ce champ.
                                  Comportement normal, pas une erreur.
      - nom_client = "Xxx"     → un nom a été renseigné mais ne correspond à
                                  aucune colonne de la couche. C'est une erreur
                                  de configuration de mapping.py, pas une
                                  absence de donnée chez le client.

    Le deuxième cas arrête le script avec un message clair : nom du champ
    fautif, couche concernée, et liste des colonnes réellement disponibles
    pour faciliter la correction immédiate de mapping.py.

    Sans cette vérification, une faute de frappe sur "nom_client" produirait
    un rapport silencieusement faux — affichant "100% non renseigné" pour un
    champ qui existe en réalité dans les données, sous un autre nom.
    """
    erreurs = []
    for cle_champ, config_champ in champs_config.items():
        nom_client = config_champ.get("nom_client")
        # nom_client = None est un choix volontaire — jamais une erreur
        if nom_client is None:
            continue
        # nom_client renseigné mais absent des colonnes → erreur de mapping
        if nom_client not in gdf.columns:
            erreurs.append((cle_champ, config_champ.get("label", cle_champ), nom_client))

    if erreurs:
        print(f"\n  [ERREUR] Configuration invalide dans mapping.py — couche {nom_couche}")
        print("  Les champs suivants ne correspondent à aucune colonne réelle :")
        for cle_champ, label, nom_client in erreurs:
            print(f"    - {cle_champ} ({label}) → \"{nom_client}\" introuvable")
        print(f"\n  Colonnes disponibles dans la couche {nom_couche} :")
        print(f"    {', '.join(list(gdf.columns))}")
        print("\n  → Corriger \"nom_client\" dans mapping.py pour chaque champ listé ci-dessus.")
        print("  → Si le champ n'existe vraiment pas chez le client, écrire nom_client: None")
        sys.exit(1)


# =============================================================================
#  LOGO SCE
# =============================================================================
def construire_logo_html():
    """
    Retourne la balise HTML du logo SCE.

    Si LOGO_SCE_PATH est renseigné dans config.py et le fichier existe :
      → Le logo est encodé en base64 et intégré directement dans le HTML.
        Cela garantit son affichage dans le PDF sans dépendance de chemin.

    Sinon :
      → Un SVG de substitution reprenant les formes du logo SCE est utilisé.
        Pour utiliser le vrai logo : renseigner LOGO_SCE_PATH dans config.py.
    """
    if LOGO_SCE_PATH and Path(LOGO_SCE_PATH).exists():
        # Lire le fichier image et l'encoder en base64
        with open(LOGO_SCE_PATH, "rb") as f:
            data = b64encode(f.read()).decode("utf-8")
        # Détecter le type MIME selon l'extension
        ext = Path(LOGO_SCE_PATH).suffix.lower().replace(".", "")
        mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
        return (
            f'<img src="data:{mime};base64,{data}" '
            f'style="width:36px;height:36px;object-fit:contain;" alt="SCE">'
        )

    # SVG de substitution (à remplacer en renseignant LOGO_SCE_PATH dans config.py)
    return '''<svg width="36" height="36" viewBox="0 0 100 100">
      <path d="M30 10 Q20 50 35 80 Q39 88 46 90 Q38 78 34 58 Q31 38 38 16Z"
            fill="#8b6f5e"/>
      <path d="M70 10 Q80 50 65 80 Q61 88 54 90 Q62 78 66 58 Q69 38 62 16Z"
            fill="#9d88c2"/>
      <path d="M50 42 Q65 52 90 65 Q76 60 60 54 Q50 50 48 48
               Q58 54 72 63 Q84 70 94 76 Q74 72 58 62 Q50 56 52 46Z"
            fill="#ed7d31"/>
    </svg>'''


# =============================================================================
#  ENCODAGE BASE64 DES IMAGES
# =============================================================================
def b64img(chemin_img):
    """
    Encode une image PNG en base64 pour l'intégrer dans le HTML.
    Nécessaire pour que xhtml2pdf et Chrome trouvent les images dans le PDF
    sans dépendre des chemins absolus du système.
    Retourne une chaîne "data:image/png;base64,..." ou chaîne vide si absent.
    """
    p = Path(chemin_img)
    if not p.exists():
        return ""   # l'image n'existe pas encore
    with open(p, "rb") as f:
        data = b64encode(f.read()).decode("utf-8")
    return f"data:image/png;base64,{data}"


# =============================================================================
#  GÉNÉRATION PDF
# =============================================================================
def generer_pdf(html_str, pdf_path):
    """
    Convertit le HTML en PDF.

    Stratégie :
      1. Tenter Chrome headless — respecte parfaitement le CSS d'impression
         (@page A4 landscape, grid, flexbox) → rendu fidèle au HTML.
      2. Si Chrome absent → fallback xhtml2pdf (moins fidèle mais fonctionne
         sans installation supplémentaire).

    Chrome headless est disponible si 'google-chrome', 'chromium-browser'
    ou 'chromium' est installé et accessible dans le PATH.
    """
    # Sauvegarder le HTML dans un fichier temporaire pour Chrome
    html_temp = pdf_path.parent / "_temp_rapport.html"
    html_temp.write_text(html_str, encoding="utf-8")

    # Chercher Chrome ou Chromium dans le PATH
    chrome_cmd = None
    for nom in ["google-chrome", "chromium-browser", "chromium",
                "chrome", "Google Chrome"]:
        if shutil.which(nom):
            chrome_cmd = nom
            break

    # Sur Windows, chercher Chrome dans les emplacements standard
    if chrome_cmd is None:
        for chemin_win in [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        ]:
            if Path(chemin_win).exists():
                chrome_cmd = chemin_win
                break

    if chrome_cmd:
        # Utiliser Chrome headless pour une conversion fidèle
        print("  Utilisation de Chrome headless pour le PDF …")
        try:
            subprocess.run([
                chrome_cmd,
                "--headless",
                "--disable-gpu",
                "--no-sandbox",
                f"--print-to-pdf={str(pdf_path)}",
                "--print-to-pdf-no-header",
                str(html_temp.resolve())
            ], check=True, capture_output=True, timeout=60)
            print(f"  PDF (Chrome) → {pdf_path}")
            html_temp.unlink(missing_ok=True)   # supprimer le fichier temporaire
            return
        except Exception as e:
            print(f"  Chrome a échoué ({e}) — bascule sur xhtml2pdf")

    # Fallback : xhtml2pdf
    print("  Utilisation de xhtml2pdf (fallback) …")
    with open(str(pdf_path), "wb") as f:
        pisa.CreatePDF(html_str, dest=f)
    print(f"  PDF (xhtml2pdf) → {pdf_path}")
    html_temp.unlink(missing_ok=True)


# =============================================================================
#  GÉNÉRATION POWERPOINT ÉDITABLE
# =============================================================================
def generer_pptx(context, cht_dir, map_dir, out_dir):
    """
    Génère un fichier PowerPoint éditable reproduisant la mise en page du rapport.
    Tous les éléments (textes, formes, tableaux) sont modifiables manuellement
    dans PowerPoint ou LibreOffice Impress.

    Utilise python-pptx (pip install python-pptx).
    Si non installé, la génération est ignorée avec un avertissement.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        # python-pptx non installé — pas bloquant
        print("  [INFO] python-pptx non installé — PowerPoint non généré.")
        print("         Pour l'activer : pip install python-pptx")
        return

    # Création de la présentation en format A4 paysage
    prs = Presentation()
    prs.slide_width  = Inches(11.69)   # largeur A4 paysage
    prs.slide_height = Inches(8.27)    # hauteur A4 paysage

    # Mise en page vide (sans zones de texte prédéfinies)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # --- Couleurs SCE ---
    PURPLE = RGBColor(0x3D, 0x2E, 0x4A)   # violet foncé
    ORANGE = RGBColor(0xED, 0x7D, 0x31)   # orange
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # blanc
    GREY   = RGBColor(0xD9, 0xCD, 0xE3)   # gris clair

    def add_rect(x, y, w, h, color, alpha=None):
        """Ajoute un rectangle plein de couleur donnée."""
        shape = slide.shapes.add_shape(
            1,   # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()   # pas de bordure
        return shape

    def add_text(x, y, w, h, texte, size=10, color=None,
                 bold=False, align=PP_ALIGN.LEFT):
        """Ajoute une zone de texte éditable."""
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = texte
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def add_image(x, y, w, chemin):
        """Ajoute une image PNG si elle existe."""
        if Path(chemin).exists():
            slide.shapes.add_picture(
                str(chemin), Inches(x), Inches(y), width=Inches(w)
            )

    # --- BANDEAU ---
    add_rect(0, 0, 11.69, 0.65, PURPLE)
    # Liseré orange sous le bandeau
    add_rect(0, 0.65, 11.69, 0.03, ORANGE)

    # Titre dans le bandeau
    add_text(0.6, 0.05, 7, 0.35,
             "Rapport de qualification des données SIG",
             size=14, color=WHITE, bold=True)
    add_text(0.6, 0.35, 7, 0.25,
             f"{context['projet_commune']} — {context['projet_nom']}",
             size=9, color=ORANGE)

    # Métadonnées à droite du bandeau
    meta = (f"Projet : {context['projet_nom']}  |  "
            f"Date : {context['projet_date']}  |  "
            f"Version : {context['projet_version']}\n"
            f"Couches : {context['nb_couches']}  |  "
            f"Objets : {context['nb_objets_total']}  |  Produit par SCE")
    add_text(7.5, 0.05, 4, 0.55, meta, size=7.5, color=GREY, align=PP_ALIGN.RIGHT)

    # --- BLOC 1 : SYNTHÈSE ---
    add_text(0.15, 0.75, 3, 0.22, "1  SYNTHÈSE EXÉCUTIVE",
             size=9, color=PURPLE, bold=True)
    add_rect(0.15, 0.97, 11.4, 0.015, ORANGE)   # liseré orange sous le titre

    # Jauges (images)
    add_image(0.15, 1.0,  2.5, str(cht_dir / "gauge_ifq.png"))
    add_image(2.9,  1.0,  2.5, str(cht_dir / "gauge_attr.png"))
    add_image(5.65, 1.0,  2.5, str(cht_dir / "gauge_topo.png"))

    # KPI anomalies
    add_rect(8.5, 1.0, 2.8, 1.4,
             RGBColor(0xFA, 0xF8, 0xFC))
    add_text(8.5, 1.05, 2.8, 0.22, "ANOMALIES DÉTECTÉES",
             size=7.5, color=RGBColor(0x6B, 0x58, 0x78), align=PP_ALIGN.CENTER)
    add_text(8.5, 1.35, 2.8, 0.6,
             str(context['total_anomalies']),
             size=28, color=RGBColor(0xC2, 0x5E, 0x0E),
             bold=True, align=PP_ALIGN.CENTER)
    add_text(8.5, 1.95, 2.8, 0.22,
             f"sur {context['nb_objets_total']} objets",
             size=7.5, color=RGBColor(0x8A, 0x7A, 0x98),
             align=PP_ALIGN.CENTER)

    # --- BLOC 2 : QUALITÉ ATTRIBUTAIRE ---
    add_text(0.15, 2.6, 5.5, 0.22, "2  QUALITÉ ATTRIBUTAIRE",
             size=9, color=PURPLE, bold=True)
    add_rect(0.15, 2.82, 5.5, 0.015, ORANGE)

    # Donut erreurs
    add_image(0.15, 2.85, 2.0, str(cht_dir / "donut_erreurs.png"))

    # Barres complétude
    add_image(0.15, 4.5, 5.5, str(cht_dir / "bar_completude.png"))

    # --- BLOC 3 : TOPOLOGIQUE ---
    add_text(6.0, 2.6, 5.5, 0.22, "3  QUALITÉ SPATIALE / TOPOLOGIQUE",
             size=9, color=PURPLE, bold=True)
    add_rect(6.0, 2.82, 5.5, 0.015, ORANGE)

    # Donut nœuds
    add_image(6.0, 2.85, 1.8, str(cht_dir / "donut_noeuds.png"))

    # Cartes
    add_image(6.0, 3.8, 5.5, str(map_dir / "map_noeuds.png"))
    add_image(6.0, 5.8, 5.5, str(map_dir / "map_troncons.png"))

    # --- BLOC 5 : ACTIONS CORRECTIVES ---
    add_text(0.15, 6.2, 5.5, 0.22, "5  ACTIONS CORRECTIVES PRIORITAIRES",
             size=9, color=PURPLE, bold=True)
    add_rect(0.15, 6.42, 5.5, 0.015, ORANGE)

    # Tableau des priorités
    y_prio = 6.5
    for i, p in enumerate(context['priorities'][:6]):
        sev_color = {
            'r': RGBColor(0xA6, 0x30, 0x30),
            'o': RGBColor(0xC2, 0x5E, 0x0E),
            'p': RGBColor(0x6B, 0x4A, 0x8A)
        }.get(p['sev'], RGBColor(0x1A, 0x1A, 0x1A))

        add_text(0.15, y_prio + i * 0.22, 0.3, 0.2,
                 str(i + 1), size=8, color=sev_color, bold=True)
        add_text(0.55, y_prio + i * 0.22, 4.5, 0.2,
                 p['desc'], size=8, color=RGBColor(0x1A, 0x1A, 0x1A))
        add_text(5.1, y_prio + i * 0.22, 0.5, 0.2,
                 str(p['count']), size=8, color=sev_color,
                 bold=True, align=PP_ALIGN.RIGHT)

    # --- FOOTER ---
    add_rect(0, 8.05, 11.69, 0.22, PURPLE)
    add_rect(0, 8.05, 11.69, 0.02, ORANGE)
    add_text(0.15, 8.07, 6, 0.18,
             f"SCE — Qualification SIG {context['projet_version']} · {context['projet_date']}",
             size=7, color=GREY)
    add_text(6, 8.07, 5.5, 0.18,
             "Document produit à partir des données transmises par le client",
             size=7, color=GREY, align=PP_ALIGN.RIGHT)

    # Sauvegarde du fichier PowerPoint
    pptx_path = out_dir / f"Rapport_qualif_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    prs.save(str(pptx_path))
    print(f"  PowerPoint → {pptx_path}")


# =============================================================================
#  PIPELINE PRINCIPAL
# =============================================================================
def main():
    args    = parse_args()
    out_dir = Path(args.out)
    cht_dir = out_dir / "charts"
    map_dir = out_dir / "maps"

    # Créer les dossiers de sortie nécessaires
    for d in [cht_dir, map_dir]:
        d.mkdir(parents=True, exist_ok=True)

    # Générer la date automatiquement en français
    try:
        locale.setlocale(locale.LC_TIME, "fr_FR.UTF-8")
    except locale.Error:
        pass   # locale française non disponible — date en anglais
    date_rapport = datetime.now().strftime("%d %B %Y")

    print("\n" + "="*50)
    print(f"  Qualification SIG — SCE V8")
    print(f"  {PROJET_NOM} — {PROJET_COMMUNE}  |  {date_rapport}")
    print("="*50)

    # -----------------------------------------------------------------------
    #  ÉTAPE 1 : Chargement des couches du client
    # -----------------------------------------------------------------------
    print("\n[1/5] Chargement …")
    gdf_troncon = charger_couche(args.troncon, "Tronçon")
    gdf_noeud   = charger_couche(args.noeud,   "Nœud")

    # Validation stricte des champs de mapping.py — arrête le script si un
    # nom_client renseigné ne correspond à aucune colonne réelle (V22).
    # Ne bloque jamais sur un champ volontairement laissé à None.
    valider_champs_mapping(gdf_troncon, CHAMPS_TRONCON, "Tronçon")
    valider_champs_mapping(gdf_noeud,   CHAMPS_NOEUD,   "Nœud")
    print("  ✓ Tous les champs renseignés dans mapping.py ont été retrouvés dans les données")

    # -----------------------------------------------------------------------
    #  ÉTAPE 2 : Qualité attributaire
    #  Utilise CHAMPS_TRONCON et CHAMPS_NOEUD depuis mapping.py.
    #  Pour un nouveau client : modifier uniquement "nom_client" dans mapping.py.
    # -----------------------------------------------------------------------
    print("\n[2/5] Contrôles attributaires …")
    qa_tr = run_quality_attributaire(gdf_troncon, CHAMPS_TRONCON, "Tronçon")
    qa_nd = run_quality_attributaire(gdf_noeud,   CHAMPS_NOEUD,   "Nœud")
    print(f"  Tronçon : complétude {qa_tr['completude_moyenne']}%  "
          f"| {qa_tr['n_objets_en_erreur']} objets avec erreur")
    print(f"  Nœud    : complétude {qa_nd['completude_moyenne']}%  "
          f"| {qa_nd['n_objets_en_erreur']} objets avec erreur")

    # -----------------------------------------------------------------------
    #  ÉTAPE 3 : Qualité topologique
    # -----------------------------------------------------------------------
    print("\n[3/5] Contrôles topologiques …")
    qt = run_quality_spatiale(gdf_troncon, gdf_noeud)
    print(f"  Nœuds non raccordés       : {qt['noeuds']['non_connectes']}")
    print(f"  Nœuds au milieu tronçon   : {qt['noeuds']['au_milieu']}")
    print(f"  Tronçons 1 extrémité libre: {qt['troncons']['un_cote_seul']}")
    print(f"  Tronçons 2 extrémités lib.: {qt['troncons']['non_delimite']}")

    # Calcul de l'indice final de qualité (IFQ)
    # Moyenne pondérée des deux scores (pondérations dans config.py)
    # arrondi_pct évite d'afficher 100% si la valeur réelle est < 100%
    score_attr = arrondi_pct(
        (qa_tr["taux_conformite"] + qa_nd["taux_conformite"]) / 2
    )
    score_topo = qt["taux_conformite_topo"]
    ifq = arrondi_pct(POIDS_ATTR * score_attr + POIDS_TOPO * score_topo)
    print(f"\n  IFQ = {ifq}%  (attributaire {score_attr}%  ·  topologique {score_topo}%)")

    # Calcul de l'indice de connaissance patrimoniale SIG (section 4.2)
    icgp = _preparer_icgp(qa_tr, qa_nd, qt)
    print(f"  ICGP = {icgp['total_score']}/40 — {icgp['niveau_global']}")

    # Synthèse transversale 4.4 — lecture pure de 4.1, 4.2, 4.3 (aucun recalcul)
    synthese_4_4 = _preparer_synthese_4_4(
        ifq, icgp, qa_tr.get("stats_classe", {}), qa_nd.get("stats_classe", {})
    )
    print(f"  4.4 : {synthese_4_4['phrase_principale']}")

    # -----------------------------------------------------------------------
    #  ÉTAPE 4 : Graphiques et cartes
    # -----------------------------------------------------------------------
    print("\n[4/5] Graphiques et cartes …")

    # Jauges demi-cercle — sans texte sous la jauge (supprimé en V5)
    jauge_demi_cercle(ifq,        couleur_score(ifq),        str(cht_dir/"gauge_ifq.png"))
    jauge_demi_cercle(score_attr, couleur_score(score_attr), str(cht_dir/"gauge_attr.png"))
    jauge_demi_cercle(score_topo, couleur_score(score_topo), str(cht_dir/"gauge_topo.png"))

    # Donut nœuds — sans légende interne (dans le HTML à côté)
    nc = qt["noeuds"]["non_connectes"]
    co = qt["noeuds"]["connectes"]
    donut(
        labels=["Raccordés", "Non raccordés"],
        valeurs=[max(co, 0.001), max(nc, 0.001)],
        couleurs=["#6B4A8A", "#ED7D31"],
        path=str(cht_dir/"donut_noeuds.png"),
        texte_centre={
            "val":   f"{qt['noeuds']['pct_connectes']}%",
            "color": "#6B4A8A",
            "sub":   "raccordés"
        }
    )

    # Donut types d'erreurs — sans légende interne
    merged = {}
    for k, v in {**qa_tr["type_errors"], **qa_nd["type_errors"]}.items():
        merged[k] = merged.get(k, 0) + v

    COULEURS_ERR = ["#ED7D31", "#6B4A8A", "#8B6F5E", "#C2A785"]
    has_errors = bool(merged)

    if has_errors:
        labs = list(merged.keys())
        vals = list(merged.values())
        cols = COULEURS_ERR[:len(labs)]
        donut(labels=labs, valeurs=vals, couleurs=cols,
              path=str(cht_dir/"donut_erreurs.png"))
        total_err = sum(vals)
        type_errors_legende = [
            (labs[i], {"pct": round(vals[i]/total_err*100), "couleur": cols[i]})
            for i in range(len(labs))
        ]
    else:
        donut(labels=["Aucune erreur"], valeurs=[1], couleurs=["#4A8050"],
              path=str(cht_dir/"donut_erreurs.png"))
        type_errors_legende = [("Aucune erreur", {"pct": 100, "couleur": "#4A8050"})]

    # Barres horizontales de complétude — groupées par couche (section 2.3)
    # Structure : liste de groupes, chaque groupe = {couche, champs[]}
    # Les champs sont triés par complétude croissante dans chaque groupe
    champs_completude_groupes = []
    for qa, nom_couche in [(qa_tr, "Tronçon"), (qa_nd, "Nœud")]:
        champs_groupe = []
        for cle, res in qa["champs_results"].items():
            label = res.get("label", cle)
            pct   = res.get("completude", 0.0)
            champs_groupe.append({
                "label":   label,
                "pct":     pct,
                "couleur": couleur_score(pct),
            })
        # Trier par complétude croissante dans chaque groupe
        champs_groupe.sort(key=lambda x: x["pct"])
        champs_completude_groupes.append({
            "couche": nom_couche,
            "champs": champs_groupe,
        })

    # Conserver aussi une liste plate pour les graphiques matplotlib
    champs_completude = []
    for groupe in champs_completude_groupes:
        champs_completude.extend(groupe["champs"])

    if champs_completude:
        barres_horizontales(
            [c["label"] for c in champs_completude],
            [c["pct"]   for c in champs_completude],
            [c["couleur"] for c in champs_completude],
            str(cht_dir/"bar_completude.png")
        )

    # Cartes topologiques (vraies données, fond OSM, paysage forcé)
    carte_noeuds_connectivite(qt["noeuds"]["gdf"],     str(map_dir/"map_noeuds.png"))
    carte_troncons_delimitation(qt["troncons"]["gdf"], str(map_dir/"map_troncons.png"))

    print("  OK")

    # -----------------------------------------------------------------------
    #  ÉTAPE 5 : Rapport HTML, PDF, PowerPoint
    # -----------------------------------------------------------------------
    print("\n[5/5] Rapport …")

    # -----------------------------------------------------------------------
    #  CONSTRUCTION DES ACTIONS CORRECTIVES — V8
    #  Deux catégories : Terrain (visite physique) / Carto (bureau)
    #  Chaque anomalie est associée à une action depuis ACTIONS_CORRECTIVES
    #  dans config.py.
    # -----------------------------------------------------------------------

    # --- Actions terrain (anomalies topologiques) ---
    actions_terrain = []
    actions_carto   = []

    # Nœuds non raccordés
    if qt["noeuds"]["non_connectes"] > 0:
        cat, action = ACTIONS_CORRECTIVES["noeud_non_raccorde"]
        actions_terrain.append({
            "desc":   "Nœuds non raccordés au réseau",
            "action": action,
            "count":  qt["noeuds"]["non_connectes"],
            "sev":    "r"
        })

    # Nœuds au milieu d'un tronçon (V8)
    if qt["noeuds"]["au_milieu"] > 0:
        cat, action = ACTIONS_CORRECTIVES["noeud_milieu_troncon"]
        actions_terrain.append({
            "desc":   "Nœuds situés au milieu d'un tronçon",
            "action": action,
            "count":  qt["noeuds"]["au_milieu"],
            "sev":    "r"
        })

    # Tronçons avec extrémité non reliée
    prob_tr = qt["troncons"]["un_cote_seul"] + qt["troncons"]["non_delimite"]
    if prob_tr > 0:
        cat, action = ACTIONS_CORRECTIVES["troncon_ext_libre"]
        actions_terrain.append({
            "desc":   "Tronçons avec extrémité non reliée à un nœud",
            "action": action,
            "count":  prob_tr,
            "sev":    "r"
        })

    # --- Actions depuis les contrôles attributaires ---
    # Mapping entre clé interne et clé ACTIONS_CORRECTIVES
    CLE_ACTION = {
        "cote_tn":         "cote_tn_manquante",
        "cote_rad":        "cote_rad_manquante",
        "diametre":        "diametre_absent",
        "hauteur":         "hauteur_absente_fosse",
        "diametre_regard": "diametre_regard_absent",
        "type_troncon":    "type_non_renseigne",
        "type_noeud":      "type_non_renseigne",
        "materiau":        "materiau_incoherent",
        "id":              "identifiant_non_conforme",
        "forme":           "caracteres_speciaux",
        "classe":          "classe_absente",         # champ classe absent
    }

    for qa in [qa_tr, qa_nd]:
        for cle, res in qa["champs_results"].items():
            label = res.get("label", cle)
            nb_non = res.get("n_invalides", 0) + res.get("n_null", 0)

            if res.get("absent"):
                # Champ absent — toujours terrain si mesure nécessaire
                cle_action = CLE_ACTION.get(cle, "type_non_renseigne")
                cat, action = ACTIONS_CORRECTIVES.get(
                    cle_action, ("Terrain", "Vérifier sur le terrain")
                )
                entree = {
                    "desc":   f"Champ absent — {label} ({qa['layer_name']})",
                    "action": action,
                    "count":  qa["n_total"],
                    "sev":    "r"
                }
            elif nb_non > 0:
                cle_action = CLE_ACTION.get(cle, "caracteres_speciaux")
                cat, action = ACTIONS_CORRECTIVES.get(
                    cle_action, ("Carto", "Vérifier dans la table attributaire")
                )
                entree = {
                    "desc":   f"Valeurs non conformes — {label} ({qa['layer_name']})",
                    "action": action,
                    "count":  nb_non,
                    "sev":    "o"
                }
            else:
                continue

            if cat == "Terrain":
                actions_terrain.append(entree)
            else:
                actions_carto.append(entree)

    # Ajouter les erreurs des contrôles supplémentaires — DÉDOUBLONNÉES
    # On regroupe par type d'anomalie pour éviter les répétitions
    # (ex: 500 incohérences TN/Radier → une seule ligne avec count=500)
    compteurs_suppl = {}   # {cle_action: count}
    for qa in [qa_tr, qa_nd]:
        for d in qa.get("details_supplementaires", []):
            p = (d.get("probleme") or "").lower()
            couche = qa["layer_name"]
            if "tn" in p and "radier" in p:
                cle = f"cote_tn_inf_rad_{couche}"
                compteurs_suppl[cle] = compteurs_suppl.get(cle, {
                    "cle_action": "cote_tn_inf_rad",
                    "desc": f"Incohérence cote TN / cote radier ({couche})",
                    "cat": "Terrain", "count": 0, "sev": "r"
                })
                compteurs_suppl[cle]["count"] += 1
            elif "non attendue" in p:
                # Information non attendue (Vision 1 — matériau/hauteur/diamètre)
                label = d.get("label", "")
                cle = f"non_attendu_{label}_{couche}"
                compteurs_suppl[cle] = compteurs_suppl.get(cle, {
                    "cle_action": "materiau_incoherent",
                    "desc": f"Information non attendue — {label} ({couche})",
                    "cat": "Carto", "count": 0, "sev": "o"
                })
                compteurs_suppl[cle]["count"] += 1
            elif "herbe" in p or "incohérent" in p:
                cle = f"materiau_incoherent_{couche}"
                compteurs_suppl[cle] = compteurs_suppl.get(cle, {
                    "cle_action": "materiau_incoherent",
                    "desc": f"Matériau incohérent avec le type ({couche})",
                    "cat": "Carto", "count": 0, "sev": "o"
                })
                compteurs_suppl[cle]["count"] += 1

    # Ajouter les anomalies groupées aux listes terrain/carto
    for cle, info in compteurs_suppl.items():
        cat_s, action_s = ACTIONS_CORRECTIVES.get(
            info["cle_action"],
            ("Carto", "Vérifier dans la table attributaire")
        )
        entree = {
            "desc":   info["desc"],
            "action": action_s,
            "count":  info["count"],
            "sev":    info["sev"]
        }
        if info["cat"] == "Terrain":
            actions_terrain.append(entree)
        else:
            actions_carto.append(entree)

    # Action corrective pour valeurs classe non renseignées
    from config import ACTIONS_CORRECTIVES as AC
    for qa in [qa_tr, qa_nd]:
        sc = qa.get("stats_classe", {})
        if sc.get("disponible"):
            # Champ présent — vérifier s'il y a des non renseignés
            non_rens = sum(
                s["valeur"] for s in sc.get("stats", [])
                if s["classe"] == "Non renseigné"
            )
            if non_rens > 0:
                cat, action = AC.get(
                    "classe_non_renseignee",
                    ("Carto", "Compléter les classes de précision dans la table attributaire")
                )
                actions_carto.append({
                    "desc":   f"Classe de précision non renseignée ({qa['layer_name']})",
                    "action": action,
                    "count":  int(non_rens),
                    "sev":    "o"
                })

    # Tri déterministe et décroissant par volume uniquement
    # L'anomalie avec le plus grand nombre d'objets apparaît en premier.
    # Pas de tri par couche — le volume prime sur tout.
    actions_terrain.sort(key=lambda x: (-x["count"], x.get("desc", "")))
    actions_carto.sort(key=lambda x:   (-x["count"], x.get("desc", "")))

    # --- Calcul du nombre d'ouvrages à visiter (terrain uniquement) ---
    # Uniquement les anomalies nécessitant une visite physique.
    # Champs carto (caractères spéciaux, identifiants) exclus.
    # Dédoublonnage : un objet avec plusieurs problèmes terrain = 1 visite.
    indices_terrain = set()

    # Nœuds topologiques
    nd_gdf = qt["noeuds"]["gdf"]
    indices_terrain.update(nd_gdf[nd_gdf["connecte"] == "non"].index.tolist())

    # Tronçons non délimités
    tr_gdf = qt["troncons"]["gdf"]
    indices_terrain.update(
        tr_gdf[tr_gdf["statut_delimitation"] != "deux_cotes"].index.tolist()
    )

    # Cotes altimétriques manquantes ou incohérentes — nœuds seulement
    for col in ["cote_tn_qualif", "cote_rad_qualif"]:
        if col in qa_nd["gdf_enrichi"].columns:
            indices_terrain.update(
                qa_nd["gdf_enrichi"][qa_nd["gdf_enrichi"][col] == "non"].index.tolist()
            )

    # Diamètre absent pour types obligatoires — tronçons
    if "diametre_qualif" in qa_tr["gdf_enrichi"].columns:
        indices_terrain.update(
            qa_tr["gdf_enrichi"][qa_tr["gdf_enrichi"]["diametre_qualif"] == "non"].index.tolist()
        )

    nb_ouvrages_a_visiter = len(indices_terrain)

    # Calcul séparé du nb de corrections bureau
    indices_bureau = set()
    CHAMPS_CARTO = ["type_troncon_qualif", "forme_qualif",
                    "materiau_qualif", "id_qualif",
                    "type_noeud_qualif"]
    for col in CHAMPS_CARTO:
        for qa in [qa_tr, qa_nd]:
            if col in qa["gdf_enrichi"].columns:
                indices_bureau.update(
                    qa["gdf_enrichi"][qa["gdf_enrichi"][col] == "non"].index.tolist()
                )
    nb_corrections_bureau = len(indices_bureau)

    # Pour compatibilité template — liste unifiée triée
    priorities = actions_terrain + actions_carto

    # --- Exemples 2.5 — UN SEUL EXEMPLE PAR TYPE DE PROBLÈME ---
    # On diversifie les exemples pour montrer des types d'erreurs différents.
    # On prend le premier exemple de chaque type de problème unique.
    # Ainsi si tous les diamètres sont non multiples de 25, on n'en montre qu'un.
    types_probleme_vus = set()   # types de problème déjà représentés
    exemples_erreurs = []

    for qa in [qa_tr, qa_nd]:
        for cle, res in qa["champs_results"].items():
            for d in res.get("details", []):
                if not d.get("valeur"):
                    continue
                # Clé de diversification : couche + champ (pas la valeur)
                # → un seul exemple par champ problématique par couche
                cle_type = (qa["layer_name"], res.get("label", cle))
                if cle_type in types_probleme_vus:
                    continue
                types_probleme_vus.add(cle_type)
                exemples_erreurs.append({
                    "couche":   qa["layer_name"],
                    "label":    res.get("label", cle),
                    "valeur":   d["valeur"],
                    "probleme": d.get("probleme", ""),
                })
                if len(exemples_erreurs) >= 6:   # max 6 exemples diversifiés
                    break
            if len(exemples_erreurs) >= 6:
                break
        if len(exemples_erreurs) >= 6:
            break

    # --- Statistiques globales ---
    nb_objets  = len(gdf_troncon) + len(gdf_noeud)
    # Anomalies détectées = objets UNIQUES ayant au moins une anomalie
    # Un objet avec plusieurs problèmes compte UNE seule fois
    # On collecte tous les indices en erreur et on dédoublonne
    indices_anomalies = set()
    # Topologiques
    nd_gdf = qt["noeuds"]["gdf"]
    indices_anomalies.update(nd_gdf[nd_gdf["connecte"] == "non"].index.tolist())
    tr_gdf = qt["troncons"]["gdf"]
    indices_anomalies.update(
        tr_gdf[tr_gdf["statut_delimitation"] != "deux_cotes"].index.tolist()
    )
    # Attributaires — objets uniques déjà calculés dans n_objets_en_erreur
    # On utilise les GeoDataFrames enrichis pour collecter les indices
    # geometrie_qualif est exclu : signalé uniquement en note 4.3,
    # pas comptabilisé comme anomalie de qualité classique (cf. discussion V19)
    for col in qa_tr["gdf_enrichi"].columns:
        if col.endswith("_qualif") and col != "geometrie_qualif":
            indices_anomalies.update(
                qa_tr["gdf_enrichi"][qa_tr["gdf_enrichi"][col] == "non"].index.tolist()
            )
    for col in qa_nd["gdf_enrichi"].columns:
        if col.endswith("_qualif") and col != "geometrie_qualif":
            indices_anomalies.update(
                qa_nd["gdf_enrichi"][qa_nd["gdf_enrichi"][col] == "non"].index.tolist()
            )
    total_anom = len(indices_anomalies)

    # --- Contexte Jinja2 ---
    context = {
        "projet_nom":          PROJET_NOM,
        "projet_commune":      PROJET_COMMUNE,
        "projet_date":         date_rapport,     # date automatique
        "projet_version":      PROJET_VERSION,
        "nb_couches":          2,
        "nb_objets_total":     nb_objets,
        "total_anomalies":     total_anom,
        "logo_html":           construire_logo_html(),
        "ifq":                 ifq,
        "score_attr":          score_attr,
        "score_topo":          score_topo,
        "icgp":                icgp,
        # Section 4.4 — Vue d'ensemble de qualification (lecture transversale)
        "synthese_4_4":        synthese_4_4,
        # Section 4.3 — Précision géographique
        "stats_classe_tr":     qa_tr.get("stats_classe", {}),
        "stats_classe_nd":     qa_nd.get("stats_classe", {}),
        "layers": [
            {"name": "Tronçon", "n_total": qa_tr["n_total"],
             "completude": qa_tr["completude_moyenne"],
             # n_objets_en_erreur = objets uniques en erreur (pas de double comptage)
             "n_invalides": qa_tr["n_objets_en_erreur"]},
            {"name": "Nœud",   "n_total": qa_nd["n_total"],
             "completude": qa_nd["completude_moyenne"],
             "n_invalides": qa_nd["n_objets_en_erreur"]},
        ],
        # Section 2.3 — barres 3 segments : conforme / invalide / non renseigné
        # Tous les champs définis dans mapping.py (CHAMPS_TRONCON/CHAMPS_NOEUD)
        # sont désormais affichés, qu'ils soient renseignés par le client ou
        # non. Un champ absent est ajouté en amont (quality_attributaire.py)
        # avec ses taux 0% conforme / 0% invalide / 100% non renseigné — donc
        # plus aucun champ n'a besoin d'être filtré ici (V20).
        "champs_conformite": [
            {
                "couche": qa["layer_name"],
                "champs": sorted([
                    {
                        "label":    res.get("label", cle),
                        "conforme": res.get("taux_conformite_champ", 0.0),
                        "invalide": res.get("taux_invalide_champ",   0.0),
                        "vide":     res.get("taux_null_champ",       0.0),
                    }
                    for cle, res in qa["champs_results"].items()
                ], key=lambda x: x["conforme"])
            }
            for qa in [qa_tr, qa_nd]
        ],
        # Section 2.4 — tableau conformité par type de tronçon et matériau
        "stats_type":       qa_tr.get("stats_type",     []),
        "stats_materiau":   qa_tr.get("stats_materiau", []),
        "has_errors":          has_errors,
        "type_errors_legende": type_errors_legende,
        "champs_completude":   champs_completude_groupes,   # groupé par couche pour section 2.3
        "exemples_erreurs":    exemples_erreurs,   # dédoublonnés
        "topo":                qt,
        "priorities":          priorities,
        "actions_terrain":     actions_terrain,
        "actions_carto":       actions_carto,
        "nb_ouvrages_visiter":  nb_ouvrages_a_visiter,
        "nb_corrections_bureau": nb_corrections_bureau,
        "nb_anomalies_topo": (qt["noeuds"]["non_connectes"]
                              + qt["troncons"]["un_cote_seul"]
                              + qt["troncons"]["non_delimite"]),
        # Images encodées en base64 pour le PDF (pas de chemin externe)
        "img_gauge_ifq":      b64img(cht_dir/"gauge_ifq.png"),
        "img_gauge_attr":     b64img(cht_dir/"gauge_attr.png"),
        "img_gauge_topo":     b64img(cht_dir/"gauge_topo.png"),
        "img_donut_noeuds":   b64img(cht_dir/"donut_noeuds.png"),
        "img_donut_erreurs":  b64img(cht_dir/"donut_erreurs.png"),
        "img_bar_completude": b64img(cht_dir/"bar_completude.png"),
        "img_map_noeuds":     b64img(map_dir/"map_noeuds.png"),
        "img_map_troncons":   b64img(map_dir/"map_troncons.png"),
    }

    # Rendu du template Jinja2
    env      = Environment(loader=FileSystemLoader(str(TPL_DIR)))
    template = env.get_template("report.html")
    html     = template.render(**context)

    # HTML de debug — ouvrir dans Chrome pour vérifier la mise en page
    html_path = out_dir / "report_debug.html"
    html_path.write_text(html, encoding="utf-8")
    print(f"  HTML debug → {html_path}")

    # PDF — Chrome headless si disponible, sinon xhtml2pdf
    pdf_path = out_dir / f"Rapport_qualif_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf"
    generer_pdf(html, pdf_path)

    # PowerPoint éditable — via viz/generate_pptx.py
    import viz.generate_pptx as _pptx_mod
    _pptx_mod.LOGO_SCE_PATH = LOGO_SCE_PATH
    generer_pptx(context, cht_dir, map_dir, out_dir)

    # --- Export couches enrichies avec champs _qualif + action_corrective ---
    # Tronçons
    gdf_tr_out = qt["troncons"]["gdf"].copy()
    for col in qa_tr["gdf_enrichi"].columns:
        if col.endswith("_qualif"):
            gdf_tr_out[col] = qa_tr["gdf_enrichi"][col]

    # Renseigner l'action corrective pour chaque tronçon
    # Priorité terrain > carto — si plusieurs problèmes, prendre le plus critique
    gdf_tr_out["action_corrective"] = ""
    gdf_tr_out["categorie_action"]  = ""

    # Tronçons non délimités → action terrain
    masque_nd = gdf_tr_out["statut_delimitation"] != "deux_cotes"
    cat_nd, act_nd = ACTIONS_CORRECTIVES["troncon_ext_libre"]
    gdf_tr_out.loc[masque_nd, "action_corrective"] = act_nd
    gdf_tr_out.loc[masque_nd, "categorie_action"]  = cat_nd

    # Champs attributaires non conformes → action selon catégorie
    for cle_q, cle_action in [
        ("diametre_qualif",  "diametre_absent"),
        ("hauteur_qualif",   "hauteur_absente_fosse"),
        ("materiau_qualif",  "materiau_incoherent"),
        ("type_troncon_qualif", "type_non_renseigne"),
    ]:
        if cle_q in gdf_tr_out.columns:
            masque = (gdf_tr_out[cle_q] == "non") & (gdf_tr_out["action_corrective"] == "")
            if masque.any() and cle_action in ACTIONS_CORRECTIVES:
                cat, act = ACTIONS_CORRECTIVES[cle_action]
                gdf_tr_out.loc[masque, "action_corrective"] = act
                gdf_tr_out.loc[masque, "categorie_action"]  = cat

    gdf_tr_out.to_file(str(out_dir/"Troncon_qualif.gpkg"), driver="GPKG")

    # Nœuds
    gdf_nd_out = qt["noeuds"]["gdf"].copy()
    for col in qa_nd["gdf_enrichi"].columns:
        if col.endswith("_qualif"):
            gdf_nd_out[col] = qa_nd["gdf_enrichi"][col]

    gdf_nd_out["action_corrective"] = ""
    gdf_nd_out["categorie_action"]  = ""

    # Nœuds non raccordés → action terrain
    masque_nc = gdf_nd_out["connecte"] == "non"
    cat_nc, act_nc = ACTIONS_CORRECTIVES["noeud_non_raccorde"]
    gdf_nd_out.loc[masque_nc, "action_corrective"] = act_nc
    gdf_nd_out.loc[masque_nc, "categorie_action"]  = cat_nc

    # Nœuds au milieu → action terrain
    if "au_milieu_troncon" in gdf_nd_out.columns:
        masque_am = (gdf_nd_out["au_milieu_troncon"] == "oui") & (gdf_nd_out["action_corrective"] == "")
        cat_am, act_am = ACTIONS_CORRECTIVES["noeud_milieu_troncon"]
        gdf_nd_out.loc[masque_am, "action_corrective"] = act_am
        gdf_nd_out.loc[masque_am, "categorie_action"]  = cat_am

    # Champs attributaires non conformes
    for cle_q, cle_action in [
        ("cote_tn_qualif",          "cote_tn_manquante"),
        ("cote_rad_qualif",         "cote_rad_manquante"),
        ("diametre_regard_qualif",  "diametre_regard_absent"),
        ("id_qualif",               "identifiant_non_conforme"),
    ]:
        if cle_q in gdf_nd_out.columns:
            masque = (gdf_nd_out[cle_q] == "non") & (gdf_nd_out["action_corrective"] == "")
            if masque.any() and cle_action in ACTIONS_CORRECTIVES:
                cat, act = ACTIONS_CORRECTIVES[cle_action]
                gdf_nd_out.loc[masque, "action_corrective"] = act
                gdf_nd_out.loc[masque, "categorie_action"]  = cat

    gdf_nd_out.to_file(str(out_dir/"Noeud_qualif.gpkg"), driver="GPKG")

    print(f"  GPKG tronçon → {out_dir}/Troncon_qualif.gpkg")
    print(f"  GPKG nœud    → {out_dir}/Noeud_qualif.gpkg")
    print(f"\n  ✓ IFQ = {ifq}%\n")


if __name__ == "__main__":
    main()
