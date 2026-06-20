# =============================================================================
#  generate_pptx.py — Génération PowerPoint éditable — V7
#
#  CORRECTION V7 :
#    - Positions Y calculées dynamiquement selon le contenu réel
#    - Tous les blocs présents : 1, 2, 3, 4, 5 + footer
#    - Labels au-dessus des jauges
#    - Textes de légende à côté des donuts (zones de texte éditables)
#    - Sous-titres 3.1 et 3.2 présents
#    - Cases colorées 3.2 présentes
#    - Bloc 4 complet : barre niveau + curseur + cases interprétation + note
#    - Deux colonnes : gauche (blocs 2) / droite (bloc 3)
#    - Bloc 4 et 5 en bas côte à côte
# =============================================================================

from pathlib import Path
from datetime import datetime


# Couleurs SCE — définies au niveau module pour être accessibles partout
C_PURPLE   = None
C_ORANGE   = None
C_ORANGE2  = None
C_WHITE    = None
C_GREY_LT  = None
C_GREY_MD  = None
C_GREEN    = None
C_RED      = None
C_PURPLE2  = None
C_BG_SOFT  = None
C_HEAD_TBL = None

LOGO_SCE_PATH = None   # renseigné depuis main.py avant appel


def _init_couleurs():
    """Initialise les couleurs RGBColor (nécessite python-pptx importé)."""
    global C_PURPLE, C_ORANGE, C_ORANGE2, C_WHITE, C_GREY_LT, C_GREY_MD
    global C_GREEN, C_RED, C_PURPLE2, C_BG_SOFT, C_HEAD_TBL
    from pptx.dml.color import RGBColor
    C_PURPLE   = RGBColor(0x3D, 0x2E, 0x4A)
    C_ORANGE   = RGBColor(0xED, 0x7D, 0x31)
    C_ORANGE2  = RGBColor(0xC2, 0x5E, 0x0E)
    C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    C_GREY_LT  = RGBColor(0xD9, 0xCD, 0xE3)
    C_GREY_MD  = RGBColor(0x6B, 0x58, 0x78)
    C_GREEN    = RGBColor(0x4A, 0x80, 0x50)
    C_RED      = RGBColor(0xA6, 0x30, 0x30)
    C_PURPLE2  = RGBColor(0x6B, 0x4A, 0x8A)
    C_BG_SOFT  = RGBColor(0xFA, 0xF8, 0xFC)
    C_HEAD_TBL = RGBColor(0xF2, 0xEC, 0xF5)


def generer_pptx(context, cht_dir, map_dir, out_dir):
    """
    Génère le PowerPoint éditable reproduisant fidèlement le HTML.
    Toutes les positions Y sont calculées dynamiquement selon le contenu.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  [INFO] python-pptx non installé — PowerPoint non généré.")
        print("         Pour l'activer : pip install python-pptx")
        return

    _init_couleurs()

    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # Format A4 paysage
    prs = Presentation()
    prs.slide_width  = Inches(11.69)
    prs.slide_height = Inches(8.27)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # =========================================================================
    #  FONCTIONS UTILITAIRES
    # =========================================================================

    def R(x, y, w, h, fill, border=None, bpt=0):
        """Rectangle plein."""
        shp = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if border and bpt > 0:
            shp.line.color.rgb = border
            shp.line.width = Pt(bpt)
        else:
            shp.line.fill.background()
        return shp

    def T(x, y, w, h, texte, size=9, color=None, bold=False,
          align=PP_ALIGN.LEFT, italic=False):
        """Zone de texte éditable."""
        if color is None:
            color = C_WHITE
        txb = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(texte)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txb

    def I(x, y, w, chemin):
        """Image PNG si le fichier existe."""
        if Path(chemin).exists():
            slide.shapes.add_picture(str(chemin), Inches(x), Inches(y),
                                     width=Inches(w))

    def titre_bloc(x, y, w, num, libelle):
        """Titre de bloc avec pastille orange et liseré."""
        R(x, y, 0.22, 0.21, C_ORANGE)
        T(x, y, 0.22, 0.21, str(num), size=9, color=C_WHITE,
          bold=True, align=PP_ALIGN.CENTER)
        T(x+0.26, y, w-0.26, 0.21, libelle.upper(),
          size=9, color=C_PURPLE, bold=True)
        R(x, y+0.23, w, 0.012, C_ORANGE)
        return y + 0.26   # retourne le Y suivant disponible

    def couleur_pct(pct):
        """Couleur selon niveau de complétude."""
        if pct >= 80: return C_GREEN
        if pct >= 50: return C_ORANGE
        return C_RED

    def sev_col(sev):
        """Couleur de criticité."""
        return {'r': C_RED, 'o': C_ORANGE2, 'p': C_PURPLE2}.get(
            sev, C_PURPLE)

    MARG  = 0.15   # marge gauche/droite générale
    COLW  = 5.55   # largeur de chaque colonne
    COLD  = MARG + COLW + 0.19   # X colonne droite

    # =========================================================================
    #  BANDEAU
    # =========================================================================
    R(0, 0, 11.69, 0.62, C_PURPLE)
    R(0, 0.62, 11.69, 0.025, C_ORANGE)

    # Logo
    if LOGO_SCE_PATH and Path(LOGO_SCE_PATH).exists():
        I(MARG, 0.1, 0.42, LOGO_SCE_PATH)
    else:
        R(MARG, 0.1, 0.42, 0.42, C_PURPLE2)
        T(MARG, 0.1, 0.42, 0.42, "SCE", size=10, color=C_WHITE,
          bold=True, align=PP_ALIGN.CENTER)

    T(0.65, 0.06, 7.0, 0.32,
      "Rapport de qualification des données SIG",
      size=13, color=C_WHITE, bold=True)
    T(0.65, 0.35, 7.0, 0.22,
      f"{context['projet_commune']} — {context['projet_nom']}",
      size=9, color=C_ORANGE)

    meta = (f"Projet : {context['projet_nom']}   |   "
            f"Date : {context['projet_date']}   |   "
            f"Version : {context['projet_version']}\n"
            f"Couches : {context['nb_couches']}   |   "
            f"Objets : {context['nb_objets_total']}   |   Produit par SCE")
    T(7.8, 0.06, 3.75, 0.52, meta, size=7.5, color=C_GREY_LT,
      align=PP_ALIGN.RIGHT)

    # =========================================================================
    #  BLOC 1 — SYNTHÈSE EXÉCUTIVE
    # =========================================================================
    B1Y = 0.72
    titre_bloc(MARG, B1Y, 11.4, 1, "Synthèse exécutive")

    KW   = 2.65    # largeur KPI
    KH   = 1.52    # hauteur KPI
    KY   = B1Y + 0.32
    KPIS = [
        ("Score global\nde qualité",  "gauge_ifq.png"),
        ("Qualité\nattributaire",     "gauge_attr.png"),
        ("Qualité\ntopologique",      "gauge_topo.png"),
        ("Anomalies\ndétectées",      None),
    ]

    for i, (label_kpi, fname) in enumerate(KPIS):
        kx = MARG + i * (KW + 0.09)
        R(kx, KY, KW, KH, C_BG_SOFT)
        # Label au-dessus de la jauge — zone de texte éditable
        T(kx, KY + 0.04, KW, 0.26, label_kpi,
          size=7.5, color=C_GREY_MD, align=PP_ALIGN.CENTER)
        if fname:
            I(kx + 0.12, KY + 0.28, KW - 0.24, str(cht_dir / fname))
        else:
            T(kx, KY + 0.42, KW, 0.58,
              str(context['total_anomalies']),
              size=26, color=C_ORANGE2, bold=True, align=PP_ALIGN.CENTER)
            T(kx, KY + 1.06, KW, 0.24,
              f"sur {context['nb_objets_total']} objets",
              size=8, color=C_GREY_MD, align=PP_ALIGN.CENTER)

    # =========================================================================
    #  COLONNE GAUCHE — BLOC 2 : QUALITÉ ATTRIBUTAIRE
    # =========================================================================
    B2Y = KY + KH + 0.12
    cy  = titre_bloc(MARG, B2Y, COLW, 2, "Qualité attributaire")

    # 2.1 Tableau
    T(MARG, cy, COLW, 0.17, "2.1  Statistiques générales par couche",
      size=8, color=C_GREY_MD, bold=True)
    cy += 0.19

    R(MARG, cy, COLW, 0.21, C_HEAD_TBL)
    for ctxt, cx_off, cw in [("Couche",0.05,1.5),("Objets",1.6,1.0),
                               ("% Complétude",2.7,1.5),("Invalides",4.3,1.2)]:
        T(MARG+cx_off, cy, cw, 0.21, ctxt, size=8, color=C_PURPLE, bold=True)
    cy += 0.22

    for layer in context['layers']:
        col_c = couleur_pct(layer['completude'])
        col_i = C_RED if layer['n_invalides']>20 else (
                C_ORANGE2 if layer['n_invalides']>0 else C_GREEN)
        T(MARG+0.05, cy, 1.5, 0.2, layer['name'],
          size=8, color=C_PURPLE)
        T(MARG+1.6,  cy, 1.0, 0.2, str(layer['n_total']),
          size=8, color=C_PURPLE)
        T(MARG+2.7,  cy, 1.5, 0.2, f"{layer['completude']}%",
          size=8, color=col_c, bold=True)
        T(MARG+4.3,  cy, 1.2, 0.2, str(layer['n_invalides']),
          size=8, color=col_i, bold=True)
        cy += 0.21

    cy += 0.05

    # 2.2 Donut + légende textuelle à côté
    T(MARG, cy, COLW, 0.17, "2.2  Répartition des types d'erreurs",
      size=8, color=C_GREY_MD, bold=True)
    cy += 0.19

    I(MARG, cy, 1.75, str(cht_dir / "donut_erreurs.png"))

    # Légende textuelle à côté du donut — zones de texte éditables
    COLS_ERR = [C_ORANGE, C_PURPLE2,
                __import__('pptx.dml.color', fromlist=['RGBColor']).RGBColor(0x8B,0x6F,0x5E),
                __import__('pptx.dml.color', fromlist=['RGBColor']).RGBColor(0xC2,0xA7,0x85)]
    leg_y = cy + 0.12
    for i, (lab, info) in enumerate(context.get('type_errors_legende', [])):
        col_leg = COLS_ERR[i % len(COLS_ERR)]
        R(MARG+1.85, leg_y + i*0.24, 0.12, 0.12, col_leg)
        T(MARG+2.02, leg_y + i*0.24, 3.35, 0.22,
          f"{lab} ({info['pct']}%)", size=8, color=C_PURPLE)

    cy += 1.55

    # 2.3 Complétude par couche
    T(MARG, cy, COLW, 0.17, "2.3  Complétude des champs clés",
      size=8, color=C_GREY_MD, bold=True)
    cy += 0.19

    for groupe in context.get('champs_completude', []):
        # Sous-titre couche — zone de texte éditable
        R(MARG, cy, COLW, 0.18, C_HEAD_TBL)
        T(MARG+0.05, cy, COLW, 0.18,
          groupe['couche'], size=8, color=C_PURPLE, bold=True)
        cy += 0.20

        for champ in groupe['champs']:
            pct = champ['pct']
            col_b = couleur_pct(pct)
            T(MARG, cy, 1.45, 0.17, champ['label'],
              size=7.5, color=C_PURPLE)
            # Fond barre
            R(MARG+1.5, cy+0.04, 3.55, 0.09,
              __import__('pptx.dml.color', fromlist=['RGBColor']).RGBColor(0xEF,0xEA,0xF2))
            # Barre colorée proportionnelle
            if pct > 0:
                R(MARG+1.5, cy+0.04, 3.55*pct/100, 0.09, col_b)
            T(MARG+5.1, cy, 0.4, 0.17,
              f"{pct}%", size=7.5, color=col_b, bold=True, align=PP_ALIGN.RIGHT)
            cy += 0.19

    cy += 0.06

    # 2.4 Exemples
    T(MARG, cy, COLW, 0.17, "2.4  Exemples de valeurs problématiques",
      size=8, color=C_GREY_MD, bold=True)
    cy += 0.19

    if context.get('exemples_erreurs'):
        R(MARG, cy, COLW, 0.20, C_HEAD_TBL)
        for ctxt, cx_off, cw in [("Couche",0.05,1.0),("Champ",1.1,1.2),
                                   ("Valeur",2.4,1.0),("Problème",3.5,2.0)]:
            T(MARG+cx_off, cy, cw, 0.20, ctxt, size=8, color=C_PURPLE, bold=True)
        cy += 0.21

        for ex in context['exemples_erreurs'][:4]:
            T(MARG+0.05, cy, 1.0, 0.18, ex['couche'], size=7.5, color=C_PURPLE)
            T(MARG+1.1,  cy, 1.2, 0.18, ex['label'],  size=7.5, color=C_PURPLE)
            T(MARG+2.4,  cy, 1.0, 0.18, ex['valeur'],  size=7.5, color=C_RED)
            T(MARG+3.5,  cy, 2.0, 0.18, ex['probleme'], size=7.5, color=C_PURPLE)
            cy += 0.19
    else:
        T(MARG, cy, COLW, 0.18,
          "✓ Aucune valeur problématique détectée",
          size=8, color=C_GREEN, bold=True)

    # =========================================================================
    #  COLONNE DROITE — BLOC 3 : QUALITÉ SPATIALE / TOPOLOGIQUE
    # =========================================================================
    dy = titre_bloc(COLD, B2Y, COLW, 3, "Qualité spatiale / topologique")

    # 3.1 Sous-titre + donut + légende
    T(COLD, dy, COLW, 0.17, "3.1  Raccordement des nœuds au réseau",
      size=8, color=C_GREY_MD, bold=True)
    dy += 0.19

    I(COLD, dy, 1.75, str(cht_dir / "donut_noeuds.png"))

    # Légende textuelle à côté du donut nœuds
    topo = context['topo']
    from pptx.dml.color import RGBColor
    for i, (lab_nd, col_nd) in enumerate([
        (f"Raccordés : {topo['noeuds']['connectes']}", C_PURPLE2),
        (f"Non raccordés : {topo['noeuds']['non_connectes']}", C_ORANGE),
    ]):
        R(COLD+1.85, dy+0.18+i*0.30, 0.12, 0.12, col_nd)
        T(COLD+2.02, dy+0.16+i*0.30, 3.35, 0.22,
          lab_nd, size=9, color=C_PURPLE, bold=True)

    dy += 1.55

    # Carte nœuds
    I(COLD, dy, COLW, str(map_dir / "map_noeuds.png"))
    dy += 1.85

    # 3.2 Sous-titre + cases délimitation
    T(COLD, dy, COLW, 0.17, "3.2  Délimitation des tronçons par leurs extrémités",
      size=8, color=C_GREY_MD, bold=True)
    dy += 0.19

    case_w = (COLW - 0.12) / 3
    cases_data = [
        (str(topo['troncons']['deux_cotes']),  "Délimités aux\ndeux extrémités",
         RGBColor(0xE8,0xF0,0xE9), C_GREEN,  RGBColor(0x4A,0x5C,0x4D)),
        (str(topo['troncons']['un_cote_seul']), "Une extrémité\nnon reliée",
         RGBColor(0xFC,0xE4,0xCF), C_ORANGE2, RGBColor(0x6B,0x4A,0x1E)),
        (str(topo['troncons']['non_delimite']), "Deux extrémités\nnon reliées",
         RGBColor(0xF5,0xDA,0xDA), C_RED,    RGBColor(0x5C,0x20,0x20)),
    ]
    for i, (val, lab, bg, col_v, col_l) in enumerate(cases_data):
        cx2 = COLD + i * (case_w + 0.06)
        R(cx2, dy, case_w, 0.82, bg)
        T(cx2, dy+0.05, case_w, 0.40, val,
          size=17, color=col_v, bold=True, align=PP_ALIGN.CENTER)
        T(cx2, dy+0.44, case_w, 0.35, lab,
          size=7.5, color=col_l, align=PP_ALIGN.CENTER)
    dy += 0.88

    # Carte tronçons
    I(COLD, dy, COLW, str(map_dir / "map_troncons.png"))

    # =========================================================================
    #  LIGNE DU BAS — BLOCS 4 ET 5
    #  Positionnés dynamiquement sous le contenu le plus long des deux colonnes
    # =========================================================================
    # Prendre le Y max des deux colonnes + marge
    B45Y = max(cy, dy + 1.85) + 0.12

    # --- BLOC 4 : NIVEAU DE QUALITÉ ---
    y4 = titre_bloc(MARG, B45Y, COLW, 4, "Niveau de qualité — interprétation")

    ifq = context['ifq']
    T(MARG, y4, COLW, 0.17,
      f"Positionnement du score global ({ifq}%)",
      size=8, color=C_GREY_MD, bold=True, align=PP_ALIGN.CENTER)
    y4 += 0.20

    # Barre tricolore
    BW = COLW - 0.10
    BX = MARG + 0.05
    BH = 0.22
    R(BX,           y4, BW*0.50, BH, C_RED)
    R(BX+BW*0.50,   y4, BW*0.30, BH, C_ORANGE)
    R(BX+BW*0.80,   y4, BW*0.20, BH, C_GREEN)
    T(BX,           y4, BW*0.50, BH, "0–50% · Faible",
      size=7.5, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(BX+BW*0.50,   y4, BW*0.30, BH, "50–80% · Moyenne",
      size=7.5, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(BX+BW*0.80,   y4, BW*0.20, BH, "80–100% · Bonne",
      size=7.5, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Curseur positionné sur le score IFQ
    cursor_x = MARG + 0.05 + BW * ifq / 100 - 0.18
    T(cursor_x, y4 + BH + 0.02, 0.4, 0.17,
      f"▲ {ifq}%", size=8, color=C_PURPLE, bold=True, align=PP_ALIGN.CENTER)
    y4 += BH + 0.22

    # Trois cases d'interprétation
    iw = (COLW - 0.16) / 3
    interp_data = [
        ("0–50% · Faible",
         "Données fortement dégradées.\nIntégration déconseillée\nsans correction.",
         RGBColor(0xF5,0xDA,0xDA), C_RED,    ifq < 50),
        ("50–80% · Moyenne",
         "Anomalies significatives.\nPlan de correction ciblé\nnécessaire.",
         RGBColor(0xFC,0xE4,0xCF), C_ORANGE2, 50 <= ifq < 80),
        ("80–100% · Bonne",
         "Données exploitables.\nAméliorations mineures\npossibles.",
         RGBColor(0xE8,0xF0,0xE9), C_GREEN,   ifq >= 80),
    ]
    for i, (tit, desc, bg, col_i, actif) in enumerate(interp_data):
        ix = MARG + 0.05 + i * (iw + 0.075)
        border = C_ORANGE if actif else None
        R(ix, y4, iw, 0.78, bg, border_color=border, bpt=1.2)
        T(ix+0.04, y4+0.04, iw-0.08, 0.20, tit,
          size=7.5, color=col_i, bold=True)
        T(ix+0.04, y4+0.25, iw-0.08, 0.50, desc,
          size=7, color=C_PURPLE)
    y4 += 0.84

    # Note métier
    score_crit = (context['score_topo']
                  if context['score_topo'] < context['score_attr']
                  else context['score_attr'])
    dim_crit   = ("topologique"
                  if context['score_topo'] < context['score_attr']
                  else "attributaire")
    note = (f"Note : Score calculé à parts égales entre qualité attributaire (50%) "
            f"et topologique (50%). La qualité {dim_crit} est la plus critique : {score_crit}%.")
    R(MARG+0.05, y4, COLW-0.10, 0.38, RGBColor(0xF5,0xEE,0xE0))
    T(MARG+0.10, y4+0.03, COLW-0.20, 0.32,
      note, size=7.5, color=RGBColor(0x5C,0x3D,0x1E))

    # --- BLOC 5 : ACTIONS CORRECTIVES ---
    y5 = titre_bloc(COLD, B45Y, COLW, 5, "Actions correctives prioritaires")

    R(COLD, y5, COLW, 0.21, C_HEAD_TBL)
    T(COLD+0.05, y5, 0.28, 0.21, "#", size=8, color=C_PURPLE, bold=True)
    T(COLD+0.38, y5, 4.55, 0.21, "Anomalie identifiée",
      size=8, color=C_PURPLE, bold=True)
    T(COLD+5.05, y5, 0.46, 0.21, "Nb objets",
      size=8, color=C_PURPLE, bold=True, align=PP_ALIGN.RIGHT)
    y5 += 0.22

    for i, p in enumerate(context.get('priorities', [])[:8]):
        col_sev = sev_col(p['sev'])
        bg_pill = (RGBColor(0xF5,0xDA,0xDA) if p['sev']=='r' else
                   RGBColor(0xFC,0xE4,0xCF) if p['sev']=='o' else
                   RGBColor(0xE8,0xDE,0xEF))
        R(COLD+0.05, y5+0.03, 0.22, 0.15, bg_pill)
        T(COLD+0.05, y5+0.03, 0.22, 0.15, str(i+1),
          size=7.5, color=col_sev, bold=True, align=PP_ALIGN.CENTER)
        T(COLD+0.32, y5, 4.60, 0.19, p['desc'],
          size=8, color=C_PURPLE)
        T(COLD+5.05, y5, 0.46, 0.19, str(p['count']),
          size=8.5, color=col_sev, bold=True, align=PP_ALIGN.RIGHT)
        # Séparateur
        R(COLD, y5+0.20, COLW, 0.004,
          RGBColor(0xF0,0xEB,0xF3))
        y5 += 0.21

    # =========================================================================
    #  FOOTER
    # =========================================================================
    R(0, 8.02, 11.69, 0.02,  C_ORANGE)
    R(0, 8.04, 11.69, 0.23,  C_PURPLE)
    T(MARG, 8.06, 6.5, 0.19,
      f"SCE — Qualification SIG {context['projet_version']} · {context['projet_date']}",
      size=7.5, color=C_GREY_LT)
    T(6.0, 8.06, 5.5, 0.19,
      "Document produit à partir des données transmises par le client",
      size=7.5, color=C_GREY_LT, align=PP_ALIGN.RIGHT)

    # =========================================================================
    #  SAUVEGARDE
    # =========================================================================
    pptx_path = (Path(out_dir) /
                 f"Rapport_qualif_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")
    prs.save(str(pptx_path))
    print(f"  PowerPoint → {pptx_path}")
